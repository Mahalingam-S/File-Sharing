import os
import sys
import base64
import urllib.request
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 1. Define raw Mermaid code for the 8 diagrams
DIAGRAMS = {
    "architecture": """graph TB
    subgraph Client Layer [Frontend - React 19 / Vite]
        UI["<b>Glassmorphism UI Components</b>"]
        CTX["<b>AuthContext - Session & State</b>"]
        ROUT["<b>Protected Routes: PrivateRoute / AdminRoute</b>"]
    end

    subgraph Network Layer
        HTTP["<b>HTTPS Requests / Axios Client</b>"]
        OAUTH["<b>Google OAuth2 Client</b>"]
    end

    subgraph Backend Layer [Express REST API Server]
        SEC["<b>Security Middleware: MongoSanitize / CORS</b>"]
        AUTH["<b>Auth Middleware: JWT & RBAC</b>"]
        MUL["<b>Upload Middleware: Multer Handler</b>"]
        CTRL["<b>Controllers: Business Logic</b>"]
        CLS["<b>Lexical Document Classifier</b>"]
    end

    subgraph Data & Storage Layer
        DB[("<b>MongoDB Database (Mongoose)</b>")]
        STOR[["<b>Local Obfuscated Disk Storage</b>"]]
    end

    UI --> CTX
    CTX --> ROUT
    ROUT --> HTTP
    HTTP --> SEC

    HTTP -.-> OAUTH
    OAUTH -.-> SEC

    SEC --> AUTH
    AUTH --> CTRL
    MUL --> CTRL
    CTRL --> CLS

    CTRL --> DB
    CTRL --> STOR""",

    "auth_flow": """sequenceDiagram
    autonumber
    participant C as React Client (Vite)
    participant A as Express API (Backend)
    participant G as Google OAuth API
    participant D as MongoDB (Users Collection)

    rect rgb(240, 244, 248)
        note right of C: Standard Authentication Path
        C->>A: POST /api/auth/login (email, password)
        A->>D: Fetch user record (include +password hash)
        D-->>A: User Schema Record
        A->>A: Validate Hash (bcrypt.compare)
        A-->>C: HTTP 200: JWT Session Token + User Metadata
    end

    rect rgb(247, 250, 252)
        note right of C: Google Hybrid OAuth Path
        C->>G: Authenticate User Session (Google Login Popup)
        G-->>C: Google Auth ID Token (Credential)
        C->>A: POST /api/auth/google (Send Google ID Token)
        A->>G: Verify Client Authenticity & ID Token Signature
        G-->>A: Token Claims (Verified Name, Email)
        A->>D: Find User by Email / Auto-Register under general role
        D-->>A: User Schema Record
        A-->>C: HTTP 200: JWT Session Token + Full User Context
    end""",

    "er_diagram": """erDiagram
    USER ||--o{ FILE : owns
    USER ||--o{ FOLDER : owns
    USER ||--o{ AUDIT_LOG : triggers
    USER ||--o{ HELP_REQUEST : submits
    FOLDER ||--o{ FILE : contains
    FOLDER ||--o{ FOLDER : "parent-child hierarchy"
    FILE ||--o{ SHARED_LINK : generates""",

    "dfd_level_0": """graph TD
    User["<b>Student / Faculty / Staff User</b>"]
    Admin["<b>System Administrator</b>"]
    Core["<b>[1.0] Campus File-Sharing Platform (Core System)</b>"]
    Google["<b>Google OAuth2 Provider</b>"]
    Smtp["<b>SMTP Email Notification Server</b>"]

    User -->|<b>1. Credentials / Google Token</b>| Core
    User -->|<b>2. Upload streams / Folder actions</b>| Core
    Core -->|<b>3. Folder trees / Categorized files</b>| User
    Core -->|<b>4. Share tokens & Drop URLs</b>| User

    Admin -->|<b>5. Account approvals / Stat triggers</b>| Core
    Core -->|<b>6. Real-time Audit Logs & Support list</b>| Admin

    Core -->|<b>7. OAuth credential verification</b>| Google
    Google -->|<b>8. Verified identity profile</b>| Core

    Core -->|<b>9. Dispatch departmental notification emails</b>| Smtp""",

    "dfd_level_1": """graph TB
    Actor["<b>Student / Faculty / Staff</b>"]
    AdmActor["<b>System Administrator</b>"]

    P1["<b>[1.0] Auth & Identity Verification</b>"]
    P2["<b>[2.0] Folder & File Management</b>"]
    P3["<b>[3.0] Anonymous Drop Box Manager</b>"]
    P4["<b>[4.0] NLP Document Classifier</b>"]
    P5["<b>[5.0] System Administration & Logging</b>"]

    D1[("<b>(D1) User Profile Store</b>")]
    D2[("<b>(D2) File & Folder Store</b>")]
    D3[("<b>(D3) Obfuscated Disk Storage (storage/)</b>")]
    D4[("<b>(D4) Immutable Audit Log Store</b>")]

    Actor -->|<b>Credentials</b>| P1
    P1 -->|<b>Fetch record / validation</b>| D1
    D1 -->|<b>User Profile</b>| P1
    P1 -->|<b>Authorized JWT Session</b>| Actor

    Actor -->|<b>Upload File & Folder creation</b>| P2
    P2 -->|<b>Save file metadata</b>| D2
    P2 -->|<b>Write obfuscated file stream</b>| D3
    D2 -->|<b>Render Directory Tree</b>| P2
    P2 -->|<b>Show workspace folders/files</b>| Actor

    Actor -->|<b>Set Dropbox config & Deadline</b>| P3
    P3 -->|<b>Write Dropbox parameters</b>| D2
    Actor -->|<b>Guest: Drop coursework file</b>| P3
    P3 -->|<b>Obfuscate & store guest file</b>| D3
    P3 -->|<b>Save file metadata to target</b>| D2

    P2 -->|<b>Send originalName & mimeType</b>| P4
    P4 -->|<b>Return lexical classification badge</b>| P2

    P2 -.->|<b>Trigger event logs</b>| P5
    P3 -.->|<b>Trigger event logs</b>| P5
    P5 -->|<b>Write immutable audit record</b>| D4
    AdmActor -->|<b>Request system stats & logs</b>| P5
    D4 -->|<b>Audit records</b>| P5
    D1 -->|<b>Registration records</b>| P5
    P5 -->|<b>Display administration feeds</b>| AdmActor""",

    "dfd_level_2": """graph TD
    User["<b>Faculty / Student Owner</b>"]

    P21["<b>[2.1] Size & Security Extension Filter</b>"]
    P22["<b>[2.2] Disk Writer & Obfuscation (UUID)</b>"]
    P23["<b>[2.3] Lexical Document Classifier</b>"]
    P24["<b>[2.4] MongoDB Metadata Broker</b>"]
    P25["<b>[2.5] Security Audit Logger</b>"]

    Disk[("<b>(D3) Disk Storage</b>")]
    DB[("<b>(D2) Metadata database</b>")]
    Audit[("<b>(D4) Audit log database</b>")]

    User -->|<b>Raw File Stream & folderId</b>| P21
    P21 -->|<b>Passes verification: size < 50MB, no execs</b>| P22
    P21 -->|<b>Fail size/ext check</b>| User

    P22 -->|<b>Write obfuscated binary stream</b>| Disk
    P22 -->|<b>Emit unique storageName & ownerId</b>| P23
    
    P23 -->|<b>Analyze name patterns & MIME</b>| P24
    P24 -->|<b>Insert complete Mongoose File Document</b>| DB
    
    P24 -->|<b>Trigger security event details</b>| P25
    P25 -->|<b>Write audit record</b>| Audit
    P24 -->|<b>HTTP 201: Upload & Classification Confirmed</b>| User""",

    "upload_flow": """graph TD
    A["<b>Client UI Request</b>"] -->|<b>File Stream Transfer</b>| B["<b>Multer Middleware Handler</b>"]
    B -->|<b>Check File Size limit: 50MB</b>| C{"<b>Size <= 50MB?</b>"}
    C -->|<b>No</b>| D["<b>Reject: HTTP 400 Payload Too Large</b>"]
    C -->|<b>Yes</b>| E["<b>Scan File Extension</b>"]
    E -->|<b>Matches: .exe, .sh, .bat, .cmd</b>| F{"<b>Is Executable? Extension Check</b>"}
    F -->|<b>Yes</b>| G["<b>Reject: HTTP 403 Security Restriction</b>"]
    F -->|<b>No</b>| H["<b>Compute Obfuscated File Name</b>"]
    H -->|<b>crypto.randomUUID</b>| I["<b>Write File stream to disk storage/</b>"]
    I --> J["<b>Invoke Auto-Classification Engine</b>"]
    J --> K["<b>Insert Mongoose File Document in DB</b>"]
    K --> L["<b>Generate Immutable Audit Log</b>"]
    L --> M["<b>Return JSON response HTTP 201 Created</b>"]""",

    "drop_folder": """sequenceDiagram
    autonumber
    participant Faculty as Faculty Owner (Dashboard)
    participant Submitter as External Submitter (Guest)
    participant API as Express Router API
    participant DB as MongoDB Schema

    Faculty->>API: POST /api/folders/share/:id (Configure Drop Box + Deadline)
    API->>DB: Set Folder: isDropFolder=true, shareToken=UUID, deadline=Date
    DB-->>API: Confirm Folder Configuration
    API-->>Faculty: Return Submission Link URL (/drop/:token)

    Note over Submitter, API: Submitter navigates to the Drop Folder Submittal Page
    Submitter->>API: GET /api/folders/drop/:token
    API->>DB: Query Folder by shareToken
    DB-->>API: Return Metadata (Name, Deadline, Owner ID)
    API->>API: Assert Date.now() < Deadline
    API-->>Submitter: Render Submittals Upload Panel (Owner Hidden, Files Hidden)

    Submitter->>API: POST /api/files (Payload: File stream + shareToken)
    API->>API: Verify token validation and deadline window
    API->>API: Apply Security Filtering & Obfuscate Filename
    API->>DB: Save File Model (ownerId=FacultyOwner, folderId=DropFolderId)
    API->>DB: Create Audit Entry
    API-->>Submitter: HTTP 201: Coursework Submission Confirmed"""
}

# Ensure output diagrams folder exists
os.makedirs("d:/File-Sharing/diagrams", exist_ok=True)

# 2. Download each diagram as PNG from mermaid.ink
print("Downloading Mermaid diagrams...")
diagram_files = {}
for name, code in DIAGRAMS.items():
    # URL safe Base64 encoding
    code_bytes = code.encode('utf-8')
    base64_str = base64.urlsafe_b64encode(code_bytes).decode('utf-8')
    url = f"https://mermaid.ink/img/{base64_str}"
    
    filepath = f"d:/File-Sharing/diagrams/{name}.png"
    print(f"Fetching: {name} -> {filepath}")
    try:
        req = urllib.request.Request(
            url, 
            headers={'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64)'}
        )
        with urllib.request.urlopen(req, timeout=20) as response:
            data = response.read()
            with open(filepath, "wb") as f:
                f.write(data)
            diagram_files[name] = filepath
            print(f"Successfully downloaded {name}.png")
    except Exception as e:
        print(f"Error downloading {name}: {e}")

# 3. Create PPTX with diagrams
print("Generating presentation...")
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG_COLOR = RGBColor(10, 17, 30)       # Sleek Dark Slate Blue
CARD_BG = RGBColor(22, 32, 54)        # Frosted glass panel representation
BORDER_COLOR = RGBColor(45, 62, 96)   # Border for glass panels
TEXT_WHITE = RGBColor(241, 245, 249)  # Main text
TEXT_MUTED = RGBColor(148, 163, 184)  # Muted body text
CYAN_ACCENT = RGBColor(6, 182, 212)   # Title accent
PURPLE_ACCENT = RGBColor(168, 85, 247)# Feature highlight accent
EMERALD_GREEN = RGBColor(16, 185, 129)# Success/Advantages badge
AMBER_YELLOW = RGBColor(245, 158, 11) # Warning/Notice badge

FONT_TITLE = "Segoe UI"
FONT_BODY = "Calibri"

def apply_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    # Add decorative line at the top
    top_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08)
    )
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = CYAN_ACCENT
    top_line.line.color.rgb = CYAN_ACCENT

def add_title(slide, text):
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = text.upper()
    p.font.name = FONT_TITLE
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT
    
    # Underline accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(1.5), Inches(0.04)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PURPLE_ACCENT
    bar.line.fill.background()

def add_glass_card(slide, left, top, width, height):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = BORDER_COLOR
    card.line.width = Pt(1.5)
    return card

def add_diagram_slide(slide_title, img_key):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    apply_bg(slide)
    add_title(slide, slide_title)
    
    # Check if image downloaded successfully
    if img_key not in diagram_files:
        error_box = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.3), Inches(2.0))
        tf = error_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"[ERROR: Diagram '{img_key}' could not be fetched from Mermaid API]"
        p.font.name = FONT_TITLE
        p.font.size = Pt(20)
        p.font.color.rgb = AMBER_YELLOW
        return
        
    img_path = diagram_files[img_key]
    
    # Bounding Box for diagram card: Left=1.0, Top=1.4, Width=11.33, Height=5.3
    box_left = Inches(1.0)
    box_top = Inches(1.4)
    box_width = Inches(11.33)
    box_height = Inches(5.3)
    
    # Add glass card background for the diagram
    add_glass_card(slide, box_left, box_top, box_width, box_height)
    
    # Calculate best fit width/height preserving aspect ratio
    try:
        with Image.open(img_path) as img:
            img_w, img_h = img.size
        
        # Target sizing in Inches
        scale_w = box_width.inches / img_w
        scale_h = box_height.inches / img_h
        scale = min(scale_w, scale_h) * 0.95 # scale down slightly for margin
        
        final_w = Inches(img_w * scale)
        final_h = Inches(img_h * scale)
        
        # Center inside bounding box
        left = box_left + (box_width - final_w) / 2
        top = box_top + (box_height - final_h) / 2
        
        # Add picture
        slide.shapes.add_picture(img_path, left, top, final_w, final_h)
    except Exception as e:
        print(f"Error drawing image {img_key} on slide: {e}")
        slide.shapes.add_picture(img_path, box_left + Inches(0.5), box_top + Inches(0.5), box_width - Inches(1.0), box_height - Inches(1.0))

slide_layout = prs.slide_layouts[6] # Blank Layout

# ==================== SLIDE 1: TITLE SLIDE ====================
slide = prs.slides.add_slide(slide_layout)
apply_bg(slide)

glow = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.2), Inches(1.2), Inches(4.2), Inches(4.2))
glow.fill.solid()
glow.fill.fore_color.rgb = RGBColor(16, 28, 54)
glow.line.color.rgb = PURPLE_ACCENT
glow.line.width = Pt(1.5)

g_txt = slide.shapes.add_textbox(Inches(8.25), Inches(2.3), Inches(4.1), Inches(2))
gtf = g_txt.text_frame
gtf.word_wrap = True
gp = gtf.paragraphs[0]
gp.text = "MERN STACK\nFULL-STACK WEB APP\nSECURE & AUDITED"
gp.font.name = FONT_TITLE
gp.font.size = Pt(20)
gp.font.bold = True
gp.font.color.rgb = CYAN_ACCENT
gp.alignment = PP_ALIGN.CENTER

title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(7.5), Inches(2.2))
tf = title_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "SECURE CAMPUS FILE SHARING & DOCUMENT MANAGEMENT SYSTEM"
p.font.name = FONT_TITLE
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = TEXT_WHITE

p2 = tf.add_paragraph()
p2.text = "Final Project Review Presentation"
p2.font.name = FONT_BODY
p2.font.size = Pt(22)
p2.font.color.rgb = CYAN_ACCENT
p2.space_before = Pt(15)

author_card = add_glass_card(slide, Inches(0.8), Inches(4.0), Inches(6.8), Inches(2.6))
author_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.1), Inches(6.4), Inches(2.4))
atf = author_box.text_frame
atf.word_wrap = True

ap1 = atf.paragraphs[0]
ap1.text = "Presented by:"
ap1.font.name = FONT_BODY
ap1.font.size = Pt(14)
ap1.font.color.rgb = TEXT_MUTED

ap2 = atf.add_paragraph()
ap2.text = "MAHALINGAM S  (Reg. No: 232MCAN0099 | Roll No: 23MCA5099)"
ap2.font.name = FONT_TITLE
ap2.font.size = Pt(16)
ap2.font.bold = True
ap2.font.color.rgb = TEXT_WHITE
ap2.space_before = Pt(5)

ap3 = atf.add_paragraph()
ap3.text = "MASTER OF COMPUTER APPLICATIONS (MCA)"
ap3.font.name = FONT_BODY
ap3.font.size = Pt(14)
ap3.font.color.rgb = TEXT_MUTED

ap4 = atf.add_paragraph()
ap4.text = "Guided by: DR.T.PARIMALAM (Associate Professor & Head)"
ap4.font.name = FONT_BODY
ap4.font.size = Pt(14)
ap4.font.bold = True
ap4.font.color.rgb = PURPLE_ACCENT
ap4.space_before = Pt(10)

ap5 = atf.add_paragraph()
ap5.text = "PG & Research Dept. of Computer Science, Nandha Arts and Science College, Erode"
ap5.font.name = FONT_BODY
ap5.font.size = Pt(12)
ap5.font.color.rgb = TEXT_MUTED

# ==================== SLIDE 2: PROJECT OVERVIEW ====================
slide = prs.slides.add_slide(slide_layout)
apply_bg(slide)
add_title(slide, "Project Overview & Key Challenges")

add_glass_card(slide, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2))
prob_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(5.2), Inches(5.0))
ptf = prob_box.text_frame
ptf.word_wrap = True

pp = ptf.paragraphs[0]
pp.text = "THE CAMPUS PROBLEM"
pp.font.name = FONT_TITLE
pp.font.size = Pt(18)
pp.font.bold = True
pp.font.color.rgb = PURPLE_ACCENT

problems = [
    ("Fragmented Sharing Channels", "Academic resources are scattered across emails, chat apps, and physical USB drives."),
    ("Security Vulnerabilities", "Risks of executable file injections, malware distribution, and NoSQL query injection."),
    ("No Access Controls (RBAC)", "Students, faculty, and administrative staff share the same paths without isolated environments."),
    ("Lack of Audit Trail", "No immutable tracking of file uploads, downloads, or shares, compromising compliance.")
]
for title, desc in problems:
    p_t = ptf.add_paragraph()
    p_t.text = f"• {title}"
    p_t.font.name = FONT_TITLE
    p_t.font.size = Pt(14)
    p_t.font.bold = True
    p_t.font.color.rgb = TEXT_WHITE
    p_t.space_before = Pt(10)
    
    p_d = ptf.add_paragraph()
    p_d.text = f"  {desc}"
    p_d.font.name = FONT_BODY
    p_d.font.size = Pt(13)
    p_d.font.color.rgb = TEXT_MUTED

add_glass_card(slide, Inches(6.8), Inches(1.5), Inches(5.6), Inches(5.2))
sol_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.2), Inches(5.0))
stf = sol_box.text_frame
stf.word_wrap = True

sp = stf.paragraphs[0]
sp.text = "THE SYSTEM SOLUTION"
sp.font.name = FONT_TITLE
sp.font.size = Pt(18)
sp.font.bold = True
sp.font.color.rgb = CYAN_ACCENT

solutions = [
    ("Isolated Workspace Drives", "Separate Department Drives and private folders for organized digital storage."),
    ("Advanced Security Protections", "Uses filename obfuscation, Multer middleware restrictions, and MongoDB sanitization."),
    ("Granular Role-Based Access", "Explicit levels of operation tailored for Students, Faculty, Staff, and Admins."),
    ("Intelligent Auto-Classification", "Integrated NLP lexical engine reads filenames to auto-categorize uploads."),
    ("Expiring Links & Drop Folders", "Supports secure tokenized share links and anonymous coursework Dropboxes.")
]
for title, desc in solutions:
    s_t = stf.add_paragraph()
    s_t.text = f"✔ {title}"
    s_t.font.name = FONT_TITLE
    s_t.font.size = Pt(14)
    s_t.font.bold = True
    s_t.font.color.rgb = EMERALD_GREEN
    s_t.space_before = Pt(8)
    
    s_d = stf.add_paragraph()
    s_d.text = f"  {desc}"
    s_d.font.name = FONT_BODY
    s_d.font.size = Pt(13)
    s_d.font.color.rgb = TEXT_MUTED

# ==================== SLIDE 3: SYSTEM ARCHITECTURE OVERVIEW ====================
slide = prs.slides.add_slide(slide_layout)
apply_bg(slide)
add_title(slide, "System Architecture Overview")

layers = [
    ("FRONTEND LAYER", "React 19 & Vite", CYAN_ACCENT, 
     ["Premium Glassmorphism UI", "Dynamic CSS Mesh Background", "Responsive Tailwind Layouts", "Axios HTTP Client", "Google OAuth Popup Component"]),
    
    ("NETWORK & ROUTING", "Express API Router", PURPLE_ACCENT, 
     ["RESTful API Architecture", "JWT Session Validation", "RBAC Custom Middleware", "Multer File Stream Handler", "MongoSanitize Guard"]),
     
    ("NLP & LOGIC ENGINE", "Lexical Classifier", AMBER_YELLOW, 
     ["Heuristic keyword parsing", "MIME fallback validation", "Tokenized sharing system", "Auto-clean expired links", "Dropbox submission handler"]),
     
    ("DATA & STORAGE", "MongoDB & Disk Store", EMERALD_GREEN, 
     ["Mongoose ODM Schema validations", "Indexes (email, shareToken)", "Local Disk Storage Obfuscation", "UUID file mapping", "Immutable Audit Logging"])
]

left_positions = [Inches(0.8), Inches(3.8), Inches(6.8), Inches(9.8)]
for i, (name, subtitle, color, points) in enumerate(layers):
    add_glass_card(slide, left_positions[i], Inches(1.6), Inches(2.7), Inches(5.1))
    
    card_box = slide.shapes.add_textbox(left_positions[i] + Inches(0.15), Inches(1.7), Inches(2.4), Inches(4.9))
    ctf = card_box.text_frame
    ctf.word_wrap = True
    
    p = ctf.paragraphs[0]
    p.text = name
    p.font.name = FONT_TITLE
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = color
    
    sub = ctf.add_paragraph()
    sub.text = subtitle
    sub.font.name = FONT_TITLE
    sub.font.size = Pt(13)
    sub.font.bold = True
    sub.font.color.rgb = TEXT_WHITE
    sub.space_before = Pt(4)
    
    spc = ctf.add_paragraph()
    spc.text = "—"
    spc.font.size = Pt(10)
    spc.font.color.rgb = TEXT_MUTED
    spc.space_before = Pt(2)
    
    for pt in points:
        li = ctf.add_paragraph()
        li.text = f"• {pt}"
        li.font.name = FONT_BODY
        li.font.size = Pt(11.5)
        li.font.color.rgb = TEXT_MUTED
        li.space_before = Pt(6)

# ==================== SLIDE 4: ARCHITECTURE DIAGRAM ====================
add_diagram_slide("High-Level Architecture Diagram", "architecture")

# ==================== SLIDE 5: AUTHENTICATION FLOW DIAGRAM ====================
add_diagram_slide("Authentication & Identity Management Flow", "auth_flow")

# ==================== SLIDE 6: CORE DATABASE SCHEMAS ====================
slide = prs.slides.add_slide(slide_layout)
apply_bg(slide)
add_title(slide, "Core Database Schema & ER Model")

add_glass_card(slide, Inches(0.8), Inches(1.5), Inches(4.5), Inches(5.2))
er_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(4.1), Inches(5.0))
etf = er_box.text_frame
etf.word_wrap = True

ep = etf.paragraphs[0]
ep.text = "DATABASE ENTITY RELATIONS"
ep.font.name = FONT_TITLE
ep.font.size = Pt(16)
ep.font.bold = True
ep.font.color.rgb = CYAN_ACCENT

er_points = [
    ("USER (1) ── (N) FILE", "A user can upload and own multiple files in the system."),
    ("USER (1) ── (N) FOLDER", "Folders are created and owned by specific user accounts."),
    ("USER (1) ── (N) AUDIT_LOG", "All user actions trigger entries in the audit logs database."),
    ("FOLDER (1) ── (N) FILE", "Virtual directories catalog files via folderId Mongoose references."),
    ("FILE (1) ── (N) SHARED_LINK", "Granular share configurations generated for file tokens."),
    ("Polymorphic Schemas", "Help tickets and audit records map dynamically to trigger entities.")
]
for title, desc in er_points:
    p_t = etf.add_paragraph()
    p_t.text = f"■ {title}"
    p_t.font.name = FONT_TITLE
    p_t.font.size = Pt(13)
    p_t.font.bold = True
    p_t.font.color.rgb = TEXT_WHITE
    p_t.space_before = Pt(8)
    
    p_d = etf.add_paragraph()
    p_d.text = f"  {desc}"
    p_d.font.name = FONT_BODY
    p_d.font.size = Pt(12)
    p_d.font.color.rgb = TEXT_MUTED

add_glass_card(slide, Inches(5.6), Inches(1.5), Inches(6.8), Inches(5.2))
dict_box = slide.shapes.add_textbox(Inches(5.8), Inches(1.6), Inches(6.4), Inches(5.0))
dtf = dict_box.text_frame
dtf.word_wrap = True

dp = dtf.paragraphs[0]
dp.text = "MONGOOSE COLLECTION BLUEPRINTS"
dp.font.name = FONT_TITLE
dp.font.size = Pt(16)
dp.font.bold = True
dp.font.color.rgb = PURPLE_ACCENT

collections = [
    ("User Model (User.js)", "Stores name, email (UK), role (student, faculty, staff, admin), department, password (select: false)."),
    ("File Model (File.js)", "Stores originalName, storageName (UUID, UK), folderId, ownerId, mimeType, sizeBytes, category, department, shareToken, downloadLimit, downloadCount."),
    ("Folder Model (Folder.js)", "Tracks name, parentId (self-ref tree), ownerId, isDropFolder, deadline (Date), shareToken, department, isPublicToDepartment."),
    ("AuditLog Model (AuditLog.js)", "Immutable logs of user actions. Tracks userId, action (upload, download, share, etc.), entityType, entityId, details, ipAddress.")
]
for cname, cdesc in collections:
    c_t = dtf.add_paragraph()
    c_t.text = f"• {cname}"
    c_t.font.name = FONT_TITLE
    c_t.font.size = Pt(13.5)
    c_t.font.bold = True
    c_t.font.color.rgb = TEXT_WHITE
    c_t.space_before = Pt(8)
    
    c_d = dtf.add_paragraph()
    c_d.text = cdesc
    c_d.font.name = FONT_BODY
    c_d.font.size = Pt(12.5)
    c_d.font.color.rgb = TEXT_MUTED

# ==================== SLIDE 7: DATABASE ER DIAGRAM ====================
add_diagram_slide("Entity Relationship (ER) Diagram", "er_diagram")

# ==================== SLIDE 8: DFD LEVEL 0 ====================
add_diagram_slide("Data Flow Diagram (Level 0 - Context)", "dfd_level_0")

# ==================== SLIDE 9: DFD LEVEL 1 ====================
add_diagram_slide("Data Flow Diagram (Level 1 - Functional)", "dfd_level_1")

# ==================== SLIDE 10: DFD LEVEL 2 ====================
add_diagram_slide("Data Flow Diagram (Level 2 - Upload & Verification)", "dfd_level_2")

# ==================== SLIDE 11: WORKFLOW: SECURE UPLOAD ====================
slide = prs.slides.add_slide(slide_layout)
apply_bg(slide)
add_title(slide, "Workflow: Secure Upload & Obfuscation Steps")

steps = [
    ("1. Stream Transfer", "Client UI sends file stream via Axios multipart FormData to the backend server."),
    ("2. Multer Filter", "Checks maximum file size (50MB) and rejects executable formats (.exe, .sh, .bat)."),
    ("3. Obfuscation", "Renames physical file to a random crypto UUID, stripping original extension on disk."),
    ("4. NLP Classify", "Passes original file name and MIME to classifier engine to determine category."),
    ("5. Write & Audit", "Saves File Model metadata in MongoDB and writes an immutable, tracking Audit Log.")
]

col_width = Inches(2.2)
spacing = Inches(0.2)
start_left = Inches(0.8)

for i, (title, desc) in enumerate(steps):
    left_pos = start_left + i * (col_width + spacing)
    add_glass_card(slide, left_pos, Inches(1.8), col_width, Inches(4.5))
    
    step_box = slide.shapes.add_textbox(left_pos + Inches(0.1), Inches(1.9), col_width - Inches(0.2), Inches(4.3))
    stf = step_box.text_frame
    stf.word_wrap = True
    
    p = stf.paragraphs[0]
    p.text = title
    p.font.name = FONT_TITLE
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT
    p.alignment = PP_ALIGN.CENTER
    
    sep = stf.add_paragraph()
    sep.text = "▼"
    sep.font.size = Pt(14)
    sep.font.color.rgb = PURPLE_ACCENT
    sep.alignment = PP_ALIGN.CENTER
    sep.space_before = Pt(5)
    
    d = stf.add_paragraph()
    d.text = desc
    d.font.name = FONT_BODY
    d.font.size = Pt(13)
    d.font.color.rgb = TEXT_MUTED
    d.space_before = Pt(10)
    d.alignment = PP_ALIGN.CENTER

# ==================== SLIDE 12: UPLOAD FLOW DIAGRAM ====================
add_diagram_slide("Secure Upload & Security Filtering Flow", "upload_flow")

# ==================== SLIDE 13: NLP CLASSIFICATION ====================
slide = prs.slides.add_slide(slide_layout)
apply_bg(slide)
add_title(slide, "NLP-Lexical Document Classification")

add_glass_card(slide, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2))
class_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(5.2), Inches(5.0))
ctf = class_box.text_frame
ctf.word_wrap = True

cp = ctf.paragraphs[0]
cp.text = "CLASSIFIER LOGIC ENGINE"
cp.font.name = FONT_TITLE
cp.font.size = Pt(18)
cp.font.bold = True
cp.font.color.rgb = PURPLE_ACCENT

rules = [
    ("Filename Pattern Mapping", "Reads original name tokens in lowercase. Priority keywords map straight to academic divisions:"),
    ("  • 'assignment', 'hw', 'homework'", "Categorized as Assignments"),
    ("  • 'lab', 'record', 'experiment'", "Categorized as Lab Records"),
    ("  • 'syllabus', 'curriculum'", "Categorized as Syllabus"),
    ("  • 'question', 'paper', 'exam'", "Categorized as Question Papers"),
    ("MIME Type Fallback Rules", "If no keywords match, maps files by standard MIME category rules:"),
    ("  • image/* → Images | video/* → Videos", "application/pdf → Documents"),
    ("  • spreadsheet / excel / csv", "Categorized as Spreadsheets")
]
for title, desc in rules:
    p_t = ctf.add_paragraph()
    p_t.text = title
    p_t.font.name = FONT_TITLE if not title.startswith("  ") else FONT_BODY
    p_t.font.size = Pt(13.5) if not title.startswith("  ") else Pt(12.5)
    p_t.font.bold = not title.startswith("  ")
    p_t.font.color.rgb = TEXT_WHITE if not title.startswith("  ") else CYAN_ACCENT
    p_t.space_before = Pt(6)
    
    p_d = ctf.add_paragraph()
    p_d.text = desc
    p_d.font.name = FONT_BODY
    p_d.font.size = Pt(12)
    p_d.font.color.rgb = TEXT_MUTED

add_glass_card(slide, Inches(6.8), Inches(1.5), Inches(5.6), Inches(5.2))
code_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.2), Inches(5.0))
cotf = code_box.text_frame
cotf.word_wrap = True

cop = cotf.paragraphs[0]
cop.text = "CLASSIFIER CODE SNAPSHOT"
cop.font.name = FONT_TITLE
cop.font.size = Pt(16)
cop.font.bold = True
cop.font.color.rgb = CYAN_ACCENT

code_text = (
    "const classifyFile = (filename, mimeType) => {\n"
    "  const name = filename.toLowerCase();\n\n"
    "  if (name.includes('assignment') || name.includes('hw'))\n"
    "    return 'Assignments';\n"
    "  if (name.includes('lab') || name.includes('record'))\n"
    "    return 'Lab Records';\n"
    "  if (name.includes('exam') || name.includes('paper'))\n"
    "    return 'Question Papers';\n"
    "  if (name.includes('syllabus'))\n"
    "    return 'Syllabus';\n\n"
    "  // Fallback Category Mapping based on MIME\n"
    "  if (mimeType.startsWith('image/')) return 'Images';\n"
    "  if (mimeType.includes('pdf')) return 'Documents';\n"
    "  if (mimeType.includes('excel')) return 'Spreadsheets';\n\n"
    "  return 'Uncategorized';\n"
    "};"
)

cp_code = cotf.add_paragraph()
cp_code.text = code_text
cp_code.font.name = "Consolas"
cp_code.font.size = Pt(11.5)
cp_code.font.color.rgb = RGBColor(186, 230, 253)
cp_code.space_before = Pt(10)

# ==================== SLIDE 14: ANONYMOUS DROP FOLDERS ====================
slide = prs.slides.add_slide(slide_layout)
apply_bg(slide)
add_title(slide, "Anonymous Coursework Drop Folders")

add_glass_card(slide, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0))
f_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(4.8))
ftf = f_box.text_frame
ftf.word_wrap = True

fp = ftf.paragraphs[0]
fp.text = "FACULTY WORKSPACE SETUP"
fp.font.name = FONT_TITLE
fp.font.size = Pt(16)
fp.font.bold = True
fp.font.color.rgb = PURPLE_ACCENT

f_points = [
    ("Dropbox Configuration", "Faculty members create a normal virtual folder and toggle the Dropbox settings, adding a submission cut-off deadline."),
    ("Tokenized Drop Links", "The API generates a secure URL hash token e.g., `/drop/:token` which is shared with students."),
    ("Security Isolation", "External/unauthenticated users can drop file streams into the box but cannot view other files inside the folder.")
]
for ft, fd in f_points:
    pt = ftf.add_paragraph()
    pt.text = f"★ {ft}"
    pt.font.name = FONT_TITLE
    pt.font.size = Pt(14)
    pt.font.bold = True
    pt.font.color.rgb = TEXT_WHITE
    pt.space_before = Pt(10)
    
    pd = ftf.add_paragraph()
    pd.text = fd
    pd.font.name = FONT_BODY
    pd.font.size = Pt(13)
    pd.font.color.rgb = TEXT_MUTED

add_glass_card(slide, Inches(6.8), Inches(1.6), Inches(5.6), Inches(5.0))
s_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.2), Inches(4.8))
stf = s_box.text_frame
stf.word_wrap = True

sp = stf.paragraphs[0]
sp.text = "SUBMITTER / GUEST OPERATION FLOW"
sp.font.name = FONT_TITLE
sp.font.size = Pt(16)
sp.font.bold = True
sp.font.color.rgb = CYAN_ACCENT

s_points = [
    ("URL Access Check", "The server verifies that current Date.now() is within the configured folder deadline range."),
    ("Clean Submission Panel", "Submitter views a simple file drop card. Other student submissions are completely hidden for privacy."),
    ("Owner Assignment & Auditing", "Uploaded file owner is set to the Faculty Creator, file is placed in target folderId, and audit entry logs the IP address.")
]
for st, sd in s_points:
    pt = stf.add_paragraph()
    pt.text = f"✔ {st}"
    pt.font.name = FONT_TITLE
    pt.font.size = Pt(14)
    pt.font.bold = True
    pt.font.color.rgb = EMERALD_GREEN
    pt.space_before = Pt(10)
    
    pd = stf.add_paragraph()
    pd.text = sd
    pd.font.name = FONT_BODY
    pd.font.size = Pt(13)
    pd.font.color.rgb = TEXT_MUTED

# ==================== SLIDE 15: DROP FOLDER SEQUENCE DIAGRAM ====================
add_diagram_slide("Anonymous Drop Folder Submissions Flow", "drop_folder")

# ==================== SLIDE 16: REST API ENDPOINTS ====================
slide = prs.slides.add_slide(slide_layout)
apply_bg(slide)
add_title(slide, "REST API Endpoints & Access Control")

routers = [
    ("AUTHENTICATION (/api/auth)", CYAN_ACCENT, [
        "POST /register — Setup accounts",
        "POST /login — Issue JWT",
        "POST /google — OAuth Auto-Register",
        "GET /me — Session profile metadata",
        "DELETE /me — Account termination"
    ]),
    ("RESOURCES (/api/files & /folders)", PURPLE_ACCENT, [
        "POST /files — Upload with checks",
        "GET /files/download/:id — Secure fetch",
        "POST /files/share/:id — Share token",
        "POST /folders — Create directory trees",
        "POST /folders/share/:id — Active Dropbox"
    ]),
    ("ADMIN & SUPPORT (/api/admin & /help)", AMBER_YELLOW, [
        "GET /admin/logs — Immutable action trail",
        "GET /admin/stats — File/user statistics",
        "PATCH /admin/users/:id/approve — RBAC gate",
        "PATCH /admin/files/:id/verify — Verify document",
        "POST /help/request — Ticket pipeline"
    ])
]

left_pos = [Inches(0.8), Inches(4.8), Inches(8.8)]
for i, (name, color, endpoints) in enumerate(routers):
    add_glass_card(slide, left_pos[i], Inches(1.6), Inches(3.7), Inches(5.1))
    
    box = slide.shapes.add_textbox(left_pos[i] + Inches(0.15), Inches(1.7), Inches(3.4), Inches(4.9))
    btf = box.text_frame
    btf.word_wrap = True
    
    p = btf.paragraphs[0]
    p.text = name
    p.font.name = FONT_TITLE
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = color
    
    spc = btf.add_paragraph()
    spc.text = "—"
    spc.font.size = Pt(10)
    spc.font.color.rgb = TEXT_MUTED
    spc.space_before = Pt(2)
    
    for ep in endpoints:
        li = btf.add_paragraph()
        li.text = ep
        li.font.name = "Consolas"
        li.font.size = Pt(11.5)
        li.font.color.rgb = TEXT_WHITE
        li.space_before = Pt(8)

# ==================== SLIDE 17: SECURITY & MITIGATIONS ====================
slide = prs.slides.add_slide(slide_layout)
apply_bg(slide)
add_title(slide, "Security Architecture & Vulnerability Mitigation")

add_glass_card(slide, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2))
v_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(5.2), Inches(5.0))
vtf = v_box.text_frame
vtf.word_wrap = True

vp = vtf.paragraphs[0]
vp.text = "MITIGATING CAMPUS SECURITY THREATS"
vp.font.name = FONT_TITLE
vp.font.size = Pt(16)
vp.font.bold = True
vp.font.color.rgb = PURPLE_ACCENT

vulnerabilities = [
    ("NoSQL Database Injection", "Attackers inject malicious query operators into email or password fields to bypass login schemes."),
    ("Directory Traversal Attacks", "Attackers pass custom paths in requests to read configurations and secret files outside target paths."),
    ("Malicious File Execution", "Uploading backdoor scripts (.php, .exe, .sh) which can run on the hosting server instance.")
]
for vt, vd in vulnerabilities:
    pt = vtf.add_paragraph()
    pt.text = f"✘ {vt}"
    pt.font.name = FONT_TITLE
    pt.font.size = Pt(13.5)
    pt.font.bold = True
    pt.font.color.rgb = AMBER_YELLOW
    pt.space_before = Pt(10)
    
    pd = vtf.add_paragraph()
    pd.text = vd
    pd.font.name = FONT_BODY
    pd.font.size = Pt(12.5)
    pd.font.color.rgb = TEXT_MUTED

add_glass_card(slide, Inches(6.8), Inches(1.5), Inches(5.6), Inches(5.2))
m_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.2), Inches(5.0))
mtf = m_box.text_frame
mtf.word_wrap = True

mp = mtf.paragraphs[0]
mp.text = "INTEGRATED SYSTEM DEFENSES"
mp.font.name = FONT_TITLE
mp.font.size = Pt(16)
mp.font.bold = True
mp.font.color.rgb = CYAN_ACCENT

defenses = [
    ("Global MongoSanitize Middleware", "Express servers run `express-mongo-sanitize` on all requests. Automatically strips prohibited keys ($ and .) used in DB injections."),
    ("Filename UUID Obfuscation", "No file retains its user-uploaded path. Obfuscated UUIDs are stored on disk in disk-storage. Original associations sit strictly behind DB security."),
    ("Multer Extension Blocklists", "Enforces checks rejecting files with blocked extensions. Mime-type signatures verify content legality prior to disk writing.")
]
for mt, md in defenses:
    pt = mtf.add_paragraph()
    pt.text = f"✔ {mt}"
    pt.font.name = FONT_TITLE
    pt.font.size = Pt(13.5)
    pt.font.bold = True
    pt.font.color.rgb = EMERALD_GREEN
    pt.space_before = Pt(10)
    
    pd = mtf.add_paragraph()
    pd.text = md
    pd.font.name = FONT_BODY
    pd.font.size = Pt(12.5)
    pd.font.color.rgb = TEXT_MUTED

# ==================== SLIDE 18: UI/UX DESIGN SYSTEM ====================
slide = prs.slides.add_slide(slide_layout)
apply_bg(slide)
add_title(slide, "UI/UX Design System & Aesthetics")

features = [
    ("GLASSMORPHISM CARDS", PURPLE_ACCENT, [
        "Floating panel representation via styling",
        "Backdrop blur rules (blur(16px))",
        "Minimalist 1px border lines",
        "Deep Slate transparent layers",
        "High contrast interactivity"
    ]),
    ("DYNAMIC MESH GRADIENT", CYAN_ACCENT, [
        "Harmonious slate HSL dark theme",
        "Animated background mesh canvas",
        "Smooth transition states on hover",
        "Inspires clean workspace interactions",
        "Eliminates static layout fatigue"
    ]),
    ("COLORED SYSTEM BADGES", EMERALD_GREEN, [
        "Assignments: Emerald Green",
        "Circulars: Amber Yellow",
        "Lab Records: Intense Cyan",
        "Syllabus: Royal Purple",
        "Lecture Notes: Sapphire Blue"
    ])
]

col_width = Inches(3.6)
spacing = Inches(0.4)
start_left = Inches(0.8)

for i, (title, color, points) in enumerate(features):
    left_pos = start_left + i * (col_width + spacing)
    add_glass_card(slide, left_pos, Inches(1.6), col_width, Inches(5.1))
    
    box = slide.shapes.add_textbox(left_pos + Inches(0.15), Inches(1.7), col_width - Inches(0.3), Inches(4.9))
    btf = box.text_frame
    btf.word_wrap = True
    
    p = btf.paragraphs[0]
    p.text = title
    p.font.name = FONT_TITLE
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = color
    
    spc = btf.add_paragraph()
    spc.text = "—"
    spc.font.size = Pt(10)
    spc.font.color.rgb = TEXT_MUTED
    spc.space_before = Pt(2)
    
    for pt in points:
        li = btf.add_paragraph()
        li.text = f"✦ {pt}"
        li.font.name = FONT_BODY
        li.font.size = Pt(13)
        li.font.color.rgb = TEXT_WHITE
        li.space_before = Pt(12)

# ==================== SLIDE 19: INSTALLATION & CONFIGURATION ====================
slide = prs.slides.add_slide(slide_layout)
apply_bg(slide)
add_title(slide, "Installation & Configuration Blueprint")

add_glass_card(slide, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2))
env_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(5.2), Inches(5.0))
etf = env_box.text_frame
etf.word_wrap = True

ep = etf.paragraphs[0]
ep.text = "ENVIRONMENT SETTINGS"
ep.font.name = FONT_TITLE
ep.font.size = Pt(16)
ep.font.bold = True
ep.font.color.rgb = CYAN_ACCENT

env_text = (
    "# Backend - backend/.env\n"
    "PORT=5000\n"
    "MONGO_URI=mongodb+srv://<db_user>:<pwd>@cluster...\n"
    "JWT_SECRET=high_entropy_secret_key_value\n"
    "EMAIL_USER=institutional-email@gmail.com\n"
    "EMAIL_PASS=gmail_app_specific_password\n\n"
    "# Frontend - frontend/.env\n"
    "VITE_API_URL=http://localhost:5000/api\n"
    "VITE_GOOGLE_CLIENT_ID=google-oauth-client-id-key\n"
)

ep_code = etf.add_paragraph()
ep_code.text = env_text
ep_code.font.name = "Consolas"
ep_code.font.size = Pt(12)
ep_code.font.color.rgb = RGBColor(186, 230, 253)
ep_code.space_before = Pt(15)

add_glass_card(slide, Inches(6.8), Inches(1.5), Inches(5.6), Inches(5.2))
cmd_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.2), Inches(5.0))
ctf = cmd_box.text_frame
ctf.word_wrap = True

cp = ctf.paragraphs[0]
cp.text = "CLI SETUP & DEPLOYMENT"
cp.font.name = FONT_TITLE
cp.font.size = Pt(16)
cp.font.bold = True
cp.font.color.rgb = PURPLE_ACCENT

commands = [
    ("Backend Package Install", "cd backend && npm install"),
    ("Frontend Package Install", "cd frontend && npm install"),
    ("Database Wipe & Seed Script", "cd backend && node wipe-and-seed.js\n(Seeds students, faculty, files, directories)"),
    ("Promote User to Admin", "node make-admin.js --email=user@domain.edu"),
    ("Launch Systems Locally", "Backend: npm run dev (Runs server on PORT 5000)\nFrontend: npm run dev (Runs Vite on PORT 5173)")
]
for title, desc in commands:
    p_t = ctf.add_paragraph()
    p_t.text = f"■ {title}"
    p_t.font.name = FONT_TITLE
    p_t.font.size = Pt(13)
    p_t.font.bold = True
    p_t.font.color.rgb = TEXT_WHITE
    p_t.space_before = Pt(8)
    
    p_d = ctf.add_paragraph()
    p_d.text = desc
    p_d.font.name = "Consolas" if "npm" in desc or "node" in desc else FONT_BODY
    p_d.font.size = Pt(11.5) if "npm" in desc or "node" in desc else Pt(12)
    p_d.font.color.rgb = CYAN_ACCENT if "npm" in desc or "node" in desc else TEXT_MUTED

# ==================== SLIDE 20: SUMMARY & QUESTIONS ====================
slide = prs.slides.add_slide(slide_layout)
apply_bg(slide)
add_title(slide, "Conclusion & Final Review Summary")

add_glass_card(slide, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2))
rc_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(5.2), Inches(5.0))
rtf = rc_box.text_frame
rtf.word_wrap = True

rp = rtf.paragraphs[0]
rp.text = "KEY PROJECT HIGHLIGHTS"
rp.font.name = FONT_TITLE
rp.font.size = Pt(18)
rp.font.bold = True
rp.font.color.rgb = CYAN_ACCENT

highlights = [
    ("Seamless Academic Workflow", "Centralized, secure space replacing scattered emails and chats."),
    ("NoSQL & Extension Defense", "Production-hardened upload sanitization blocking backdoors."),
    ("Heuristic Categorization", "NLP classifier parses files to auto-organize academic directories."),
    ("Google OAuth Integration", "Hybrid identity validation combining standard login with Google accounts."),
    ("Expiring Link Audits", "Automatic cleanups of expired links and real-time audit logging.")
]
for title, desc in highlights:
    pt = rtf.add_paragraph()
    pt.text = f"✔ {title}"
    pt.font.name = FONT_TITLE
    pt.font.size = Pt(14)
    pt.font.bold = True
    pt.font.color.rgb = EMERALD_GREEN
    pt.space_before = Pt(10)
    
    pd = rtf.add_paragraph()
    pd.text = desc
    pd.font.name = FONT_BODY
    pd.font.size = Pt(13)
    pd.font.color.rgb = TEXT_MUTED
    
add_glass_card(slide, Inches(6.8), Inches(1.5), Inches(5.6), Inches(5.2))
qa_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.2), Inches(5.0))
qtf = qa_box.text_frame
qtf.word_wrap = True

qp = qtf.paragraphs[0]
qp.text = "THANK YOU & QUESTIONS"
qp.font.name = FONT_TITLE
qp.font.size = Pt(24)
qp.font.bold = True
qp.font.color.rgb = PURPLE_ACCENT
qp.alignment = PP_ALIGN.CENTER
qp.space_before = Pt(20)

qp2 = qtf.add_paragraph()
qp2.text = "SECURE CAMPUS FILE SHARING & DOCUMENT MANAGEMENT SYSTEM"
qp2.font.name = FONT_TITLE
qp2.font.size = Pt(16)
qp2.font.bold = True
qp2.font.color.rgb = TEXT_WHITE
qp2.alignment = PP_ALIGN.CENTER
qp2.space_before = Pt(40)

qp3 = qtf.add_paragraph()
qp3.text = "PG & Research Department of Computer Science"
qp3.font.name = FONT_BODY
qp3.font.size = Pt(14)
qp3.font.color.rgb = TEXT_MUTED
qp3.alignment = PP_ALIGN.CENTER
qp3.space_before = Pt(10)

qp4 = qtf.add_paragraph()
qp4.text = "Nandha Arts and Science College, Erode"
qp4.font.name = FONT_BODY
qp4.font.size = Pt(14)
qp4.font.color.rgb = TEXT_MUTED
qp4.alignment = PP_ALIGN.CENTER

qp5 = qtf.add_paragraph()
qp5.text = "Open to Queries & Feedback"
qp5.font.name = FONT_TITLE
qp5.font.size = Pt(18)
qp5.font.bold = True
qp5.font.color.rgb = CYAN_ACCENT
qp5.alignment = PP_ALIGN.CENTER
qp5.space_before = Pt(40)

# Save presentation
output_filename = "d:/File-Sharing/Secure_Campus_File_Sharing_Presentation.pptx"
prs.save(output_filename)
print(f"Presentation with diagrams saved successfully to: {output_filename}")
