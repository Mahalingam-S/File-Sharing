import os
import sys
import base64
import urllib.request
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 1. Define raw Mermaid code for the 3 key diagrams
DIAGRAMS = {
    "architecture": """graph TB
    subgraph Client Layer [Frontend - React 19 / Vite]
        UI["<b>Glassmorphism UI Components</b>"]
        CTX["<b>AuthContext - Session & State</b>"]
        ROUT["<b>Protected Routes: PrivateRoute / AdminRoute</b>"]
    end

    subgraph Network Layer
        HTTP["<b>HTTPS Requests / Axios Client</b>"]
        OAUTH["<b>Google OAuth2 Client</b>"]
    end

    subgraph Backend Layer [Express REST API Server]
        SEC["<b>Security Middleware: MongoSanitize / CORS</b>"]
        AUTH["<b>Auth Middleware: JWT & RBAC</b>"]
        MUL["<b>Upload Middleware: Multer Handler</b>"]
        CTRL["<b>Controllers: Business Logic</b>"]
        CLS["<b>Lexical Document Classifier</b>"]
    end

    subgraph Data & Storage Layer
        DB[("<b>MongoDB Database (Mongoose)</b>")]
        STOR[["<b>Local Obfuscated Disk Storage</b>"]]
    end

    UI --> CTX
    CTX --> ROUT
    ROUT --> HTTP
    HTTP --> SEC

    HTTP -.-> OAUTH
    OAUTH -.-> SEC

    SEC --> AUTH
    AUTH --> CTRL
    MUL --> CTRL
    CTRL --> CLS

    CTRL --> DB
    CTRL --> STOR""",

    "dfd_level_0": """graph TD
    User["<b>Student / Faculty / Staff User</b>"]
    Admin["<b>System Administrator</b>"]
    Core["<b>[1.0] Campus File-Sharing Platform (Core System)</b>"]
    Google["<b>Google OAuth2 Provider</b>"]
    Smtp["<b>SMTP Email Notification Server</b>"]

    User -->|<b>1. Credentials / Google Token</b>| Core
    User -->|<b>2. Upload streams / Folder actions</b>| Core
    Core -->|<b>3. Folder trees / Categorized files</b>| User
    Core -->|<b>4. Share tokens & Drop URLs</b>| User

    Admin -->|<b>5. Account approvals / Stat triggers</b>| Core
    Core -->|<b>6. Real-time Audit Logs & Support list</b>| Admin

    Core -->|<b>7. OAuth credential verification</b>| Google
    Google -->|<b>8. Verified identity profile</b>| Core

    Core -->|<b>9. Dispatch departmental notification emails</b>| Smtp""",

    "upload_flow": """graph TD
    A["<b>Client UI Request</b>"] -->|<b>File Stream Transfer</b>| B["<b>Multer Middleware Handler</b>"]
    B -->|<b>Check File Size limit: 50MB</b>| C{"<b>Size <= 50MB?</b>"}
    C -->|<b>No</b>| D["<b>Reject: HTTP 400 Payload Too Large</b>"]
    C -->|<b>Yes</b>| E["<b>Scan File Extension</b>"]
    E -->|<b>Matches: .exe, .sh, .bat, .cmd</b>| F{"<b>Is Executable? Extension Check</b>"}
    F -->|<b>Yes</b>| G["<b>Reject: HTTP 403 Security Restriction</b>"]
    F -->|<b>No</b>| H["<b>Compute Obfuscated File Name</b>"]
    H -->|<b>crypto.randomUUID</b>| I["<b>Write File stream to disk storage/</b>"]
    I --> J["<b>Invoke Auto-Classification Engine</b>"]
    J --> K["<b>Insert Mongoose File Document in DB</b>"]
    K --> L["<b>Generate Immutable Audit Log</b>"]
    L --> M["<b>Return JSON response HTTP 201 Created</b>"]"""
}

# Ensure output diagrams folder exists
os.makedirs("d:/File-Sharing/diagrams", exist_ok=True)

# 2. Download each diagram as PNG from mermaid.ink
print("Downloading Mermaid diagrams...")
diagram_files = {}
for name, code in DIAGRAMS.items():
    # URL safe Base64 encoding
    code_bytes = code.encode('utf-8')
    base64_str = base64.urlsafe_b64encode(code_bytes).decode('utf-8')
    url = f"https://mermaid.ink/img/{base64_str}"
    
    filepath = f"d:/File-Sharing/diagrams/{name}.png"
    print(f"Fetching: {name} -> {filepath}")
    try:
        req = urllib.request.Request(
            url, 
            headers={'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64)'}
        )
        with urllib.request.urlopen(req, timeout=20) as response:
            data = response.read()
            with open(filepath, "wb") as f:
                f.write(data)
            diagram_files[name] = filepath
            print(f"Successfully downloaded {name}.png")
    except Exception as e:
        print(f"Error downloading {name}: {e}")

# 3. Create PPTX
print("Generating presentation...")
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors (Aesthetic Slate Theme)
BG_COLOR = RGBColor(10, 17, 30)       # Sleek Dark Slate Blue
CARD_BG = RGBColor(22, 32, 54)        # Frosted glass panel representation
BORDER_COLOR = RGBColor(45, 62, 96)   # Border for glass panels
TEXT_WHITE = RGBColor(241, 245, 249)  # Main text
TEXT_MUTED = RGBColor(148, 163, 184)  # Muted body text
CYAN_ACCENT = RGBColor(6, 182, 212)   # Title accent
PURPLE_ACCENT = RGBColor(168, 85, 247)# Feature highlight accent
EMERALD_GREEN = RGBColor(16, 185, 129)# Success/Advantages badge
AMBER_YELLOW = RGBColor(245, 158, 11) # Warning/Notice badge

FONT_TITLE = "Segoe UI"
FONT_BODY = "Calibri"

def apply_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    # Add decorative line at the top
    top_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08)
    )
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = CYAN_ACCENT
    top_line.line.color.rgb = CYAN_ACCENT

def add_title(slide, text):
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = text.upper()
    p.font.name = FONT_TITLE
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT
    
    # Underline accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(1.5), Inches(0.04)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PURPLE_ACCENT
    bar.line.fill.background()

def add_glass_card(slide, left, top, width, height):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = BORDER_COLOR
    card.line.width = Pt(1.5)
    return card

def add_diagram_slide(slide_title, img_key):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    apply_bg(slide)
    add_title(slide, slide_title)
    
    if img_key not in diagram_files:
        error_box = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.3), Inches(2.0))
        tf = error_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"[ERROR: Diagram '{img_key}' could not be fetched from Mermaid API]"
        p.font.name = FONT_TITLE
        p.font.size = Pt(20)
        p.font.color.rgb = AMBER_YELLOW
        return
        
    img_path = diagram_files[img_key]
    
    # Bounding Box for diagram card
    box_left = Inches(1.0)
    box_top = Inches(1.4)
    box_width = Inches(11.33)
    box_height = Inches(5.3)
    
    add_glass_card(slide, box_left, box_top, box_width, box_height)
    
    # Aspect Ratio preservation
    try:
        with Image.open(img_path) as img:
            img_w, img_h = img.size
        
        scale_w = box_width.inches / img_w
        scale_h = box_height.inches / img_h
        scale = min(scale_w, scale_h) * 0.95 # slight margin
        
        final_w = Inches(img_w * scale)
        final_h = Inches(img_h * scale)
        
        left = box_left + (box_width - final_w) / 2
        top = box_top + (box_height - final_h) / 2
        
        slide.shapes.add_picture(img_path, left, top, final_w, final_h)
    except Exception as e:
        print(f"Error drawing image {img_key} on slide: {e}")
        slide.shapes.add_picture(img_path, box_left + Inches(0.5), box_top + Inches(0.5), box_width - Inches(1.0), box_height - Inches(1.0))

slide_layout = prs.slide_layouts[6] # Blank Layout

# ==================== SLIDE 1: TITLE SLIDE ====================
slide = prs.slides.add_slide(slide_layout)
apply_bg(slide)

# Decorative Glowing shape in center-right
glow = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.2), Inches(1.2), Inches(4.2), Inches(4.2))
glow.fill.solid()
glow.fill.fore_color.rgb = RGBColor(16, 28, 54)
glow.line.color.rgb = PURPLE_ACCENT
glow.line.width = Pt(1.5)

g_txt = slide.shapes.add_textbox(Inches(8.25), Inches(2.3), Inches(4.1), Inches(2))
gtf = g_txt.text_frame
gtf.word_wrap = True
gp = gtf.paragraphs[0]
gp.text = "SMART SHARING\nSECURED STORAGE\nACADEMIC WORKSPACE"
gp.font.name = FONT_TITLE
gp.font.size = Pt(20)
gp.font.bold = True
gp.font.color.rgb = CYAN_ACCENT
gp.alignment = PP_ALIGN.CENTER

title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(7.2), Inches(2.4))
tf = title_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "SECURE CAMPUS FILE SHARING & DOCUMENT MANAGEMENT SYSTEM"
p.font.name = FONT_TITLE
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = TEXT_WHITE

p2 = tf.add_paragraph()
p2.text = "Final Project Review Presentation"
p2.font.name = FONT_BODY
p2.font.size = Pt(20)
p2.font.color.rgb = CYAN_ACCENT
p2.space_before = Pt(10)

# Author & Guide Block Card
author_card = add_glass_card(slide, Inches(0.8), Inches(3.9), Inches(7.0), Inches(2.9))
author_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.0), Inches(6.6), Inches(2.7))
atf = author_box.text_frame
atf.word_wrap = True

ap1 = atf.paragraphs[0]
ap1.text = "Presented by:"
ap1.font.name = FONT_BODY
ap1.font.size = Pt(13)
ap1.font.color.rgb = TEXT_MUTED

ap2 = atf.add_paragraph()
ap2.text = "MAHALINGAM S  (Reg. No: 232MCAN0099 | Roll No: 23MCA5099)"
ap2.font.name = FONT_TITLE
ap2.font.size = Pt(14)
ap2.font.bold = True
ap2.font.color.rgb = TEXT_WHITE
ap2.space_before = Pt(4)

ap3 = atf.add_paragraph()
ap3.text = "MASTER OF COMPUTER APPLICATIONS (MCA)"
ap3.font.name = FONT_BODY
ap3.font.size = Pt(13)
ap3.font.color.rgb = TEXT_MUTED

ap4 = atf.add_paragraph()
ap4.text = "Guided by:"
ap4.font.name = FONT_BODY
ap4.font.size = Pt(13)
ap4.font.color.rgb = TEXT_MUTED
ap4.space_before = Pt(8)

ap5 = atf.add_paragraph()
ap5.text = "Dr. K. KARTHIKEYAN, MSc., MPhil., PhD., SET.,"
ap5.font.name = FONT_TITLE
ap5.font.size = Pt(14)
ap5.font.bold = True
ap5.font.color.rgb = PURPLE_ACCENT

ap6 = atf.add_paragraph()
ap6.text = "Assistant Professor of Computer Science\nPuratchi Thalaivi Amma Government Arts and Science College, Palladam - 641664."
ap6.font.name = FONT_BODY
ap6.font.size = Pt(12)
ap6.font.color.rgb = TEXT_MUTED

# ==================== SLIDE 2: THE PROBLEM & THE VISION ====================
slide = prs.slides.add_slide(slide_layout)
apply_bg(slide)
add_title(slide, "The Challenge & The Vision")

# Left Column: Problems
add_glass_card(slide, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2))
prob_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(5.2), Inches(5.0))
ptf = prob_box.text_frame
ptf.word_wrap = True

pp = ptf.paragraphs[0]
pp.text = "THE PROBLEM"
pp.font.name = FONT_TITLE
pp.font.size = Pt(18)
pp.font.bold = True
pp.font.color.rgb = PURPLE_ACCENT

problems = [
    ("Scattered Communication Channels", "Academic materials are distributed over unmanaged chats, personal drives, and email chains."),
    ("Malicious File Transmission", "Insecure systems easily spread software bugs, malware, and malicious script formats."),
    ("No Centralized Organization", "Assignments, question papers, and syllabus copies sit unorganized in single locations."),
    ("Missing Tracking & Audit Logs", "Staff cannot verify who uploaded, modified, or downloaded sensitive official documents.")
]
for title, desc in problems:
    p_t = ptf.add_paragraph()
    p_t.text = f"• {title}"
    p_t.font.name = FONT_TITLE
    p_t.font.size = Pt(14)
    p_t.font.bold = True
    p_t.font.color.rgb = TEXT_WHITE
    p_t.space_before = Pt(8)
    
    p_d = ptf.add_paragraph()
    p_d.text = f"  {desc}"
    p_d.font.name = FONT_BODY
    p_d.font.size = Pt(12.5)
    p_d.font.color.rgb = TEXT_MUTED

# Right Column: The Vision
add_glass_card(slide, Inches(6.8), Inches(1.5), Inches(5.6), Inches(5.2))
sol_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.2), Inches(5.0))
stf = sol_box.text_frame
stf.word_wrap = True

sp = stf.paragraphs[0]
sp.text = "THE VISION"
sp.font.name = FONT_TITLE
sp.font.size = Pt(18)
sp.font.bold = True
sp.font.color.rgb = CYAN_ACCENT

solutions = [
    ("Unified Academic Drive", "A central repository for isolated department drives and virtual folders."),
    ("Safe & Sandbox Storage", "Filename obfuscation and check filters prevent unauthorized backend execution."),
    ("Auto-Sorting Index Engine", "An NLP-based file categorizer sorts files automatically on upload (e.g. Lab Records, Syllabus)."),
    ("Accountable Action Trail", "Maintains real-time history logging of every single file interaction for safety.")
]
for title, desc in solutions:
    s_t = stf.add_paragraph()
    s_t.text = f"✔ {title}"
    s_t.font.name = FONT_TITLE
    s_t.font.size = Pt(14)
    s_t.font.bold = True
    s_t.font.color.rgb = EMERALD_GREEN
    s_t.space_before = Pt(8)
    
    s_d = stf.add_paragraph()
    s_d.text = f"  {desc}"
    s_d.font.name = FONT_BODY
    s_d.font.size = Pt(12.5)
    s_d.font.color.rgb = TEXT_MUTED

# ==================== SLIDE 3: CORE SYSTEM FEATURES ====================
slide = prs.slides.add_slide(slide_layout)
apply_bg(slide)
add_title(slide, "Core Platform Features")

features = [
    ("ISOLATED WORKSPACES", CYAN_ACCENT, [
        "Department-level virtual drives",
        "Private storage folders for users",
        "Hierarchical nested folders",
        "Clean, modern UI layout"
    ]),
    ("SMART AUTO-SORTING", PURPLE_ACCENT, [
        "Parses upload filename patterns",
        "Recognizes keywords automatically",
        "Sorts into: Assignments, Lab Records, Syllabus, Lecture Notes",
        "Reduces manual directory sorting"
    ]),
    ("ROLE-BASED VISIBILITY", EMERALD_GREEN, [
        "Student dashboard for coursework",
        "Faculty control panel for classes",
        "Admin control system dashboard",
        "Restricted access permissions"
    ])
]

left_positions = [Inches(0.8), Inches(4.8), Inches(8.8)]
for i, (name, color, points) in enumerate(features):
    add_glass_card(slide, left_positions[i], Inches(1.6), Inches(3.7), Inches(5.1))
    
    card_box = slide.shapes.add_textbox(left_positions[i] + Inches(0.15), Inches(1.7), Inches(3.4), Inches(4.9))
    ctf = card_box.text_frame
    ctf.word_wrap = True
    
    p = ctf.paragraphs[0]
    p.text = name
    p.font.name = FONT_TITLE
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = color
    
    # Line divider
    spc = ctf.add_paragraph()
    spc.text = "—"
    spc.font.size = Pt(10)
    spc.font.color.rgb = TEXT_MUTED
    spc.space_before = Pt(4)
    
    for pt in points:
        li = ctf.add_paragraph()
        li.text = f"✦ {pt}"
        li.font.name = FONT_BODY
        li.font.size = Pt(13)
        li.font.color.rgb = TEXT_WHITE
        li.space_before = Pt(12)

# ==================== SLIDE 4: HIGH-LEVEL ARCHITECTURE DIAGRAM ====================
add_diagram_slide("High-Level System Architecture", "architecture")

# ==================== SLIDE 5: DFD LEVEL 0 (CONTEXT DIAGRAM) ====================
add_diagram_slide("Data Flow Diagram (Level 0 - Context)", "dfd_level_0")

# ==================== SLIDE 6: SECURE UPLOAD FLOWCHART ====================
add_diagram_slide("Secure Upload & Obfuscation Flowchart", "upload_flow")

# ==================== SLIDE 7: SPECIAL ACADEMIC TOOLS ====================
slide = prs.slides.add_slide(slide_layout)
apply_bg(slide)
add_title(slide, "Special Academic Collaboration Tools")

# Left Column: Anonymous Dropboxes
add_glass_card(slide, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2))
drop_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(5.2), Inches(5.0))
dtf = drop_box.text_frame
dtf.word_wrap = True

dp = dtf.paragraphs[0]
dp.text = "ANONYMOUS COURSEWORK DROPBOXES"
dp.font.name = FONT_TITLE
dp.font.size = Pt(18)
dp.font.bold = True
dp.font.color.rgb = CYAN_ACCENT

drop_points = [
    ("Faculty Setup Control", "Faculty creates a dropbox folder, setting a secure cutoff submission deadline."),
    ("Secure Upload-Only Link", "A custom hash URL (/drop/:token) is given to students. They drop coursework directly."),
    ("Student-to-Student Isolation", "Students cannot view or read other submissions in the folder, ensuring total grading privacy."),
    ("Automated Cutoff Enforcement", "System strictly blocks uploads when the deadline is crossed.")
]
for title, desc in drop_points:
    p_t = dtf.add_paragraph()
    p_t.text = f"★ {title}"
    p_t.font.name = FONT_TITLE
    p_t.font.size = Pt(14)
    p_t.font.bold = True
    p_t.font.color.rgb = TEXT_WHITE
    p_t.space_before = Pt(8)
    
    p_d = dtf.add_paragraph()
    p_d.text = f"  {desc}"
    p_d.font.name = FONT_BODY
    p_d.font.size = Pt(12.5)
    p_d.font.color.rgb = TEXT_MUTED

# Right Column: Controlled Share Links
add_glass_card(slide, Inches(6.8), Inches(1.5), Inches(5.6), Inches(5.2))
share_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.2), Inches(5.0))
stf = share_box.text_frame
stf.word_wrap = True

sp = stf.paragraphs[0]
sp.text = "CONTROLLED FILE SHARING"
sp.font.name = FONT_TITLE
sp.font.size = Pt(18)
sp.font.bold = True
sp.font.color.rgb = PURPLE_ACCENT

share_points = [
    ("Temporary Tokenized Links", "Files share via secure hash keys, bypassing direct path access of standard servers."),
    ("Time-Limited Expiry", "Configure share validity windows (e.g., automatically expire in 2, 12, or 24 hours)."),
    ("Download Quantity Limiters", "Links self-destruct after reaching a download count threshold (e.g., single download link)."),
    ("Instant Revocation Dashboard", "File owners can cancel any shared link instantly from their control panel.")
]
for title, desc in share_points:
    s_t = stf.add_paragraph()
    s_t.text = f"✔ {title}"
    s_t.font.name = FONT_TITLE
    s_t.font.size = Pt(14)
    s_t.font.bold = True
    s_t.font.color.rgb = EMERALD_GREEN
    s_t.space_before = Pt(8)
    
    s_d = stf.add_paragraph()
    s_d.text = f"  {desc}"
    s_d.font.name = FONT_BODY
    s_d.font.size = Pt(12.5)
    s_d.font.color.rgb = TEXT_MUTED

# ==================== SLIDE 8: SECURITY & ACCOUNTABILITY MADE SIMPLE ====================
slide = prs.slides.add_slide(slide_layout)
apply_bg(slide)
add_title(slide, "Security & Accountability")

sec_features = [
    ("FILE SHIELDING", CYAN_ACCENT, [
        "Obfuscates file filenames to UUIDs on disk",
        "Strips extension paths to stop traversal",
        "Rejects executable formats (.exe, .sh)",
        "Enforces strict upload size limit (50MB)"
    ]),
    ("VERIFICATION BADGES", PURPLE_ACCENT, [
        "Official files vetted by administrators",
        "Special 'Verified' badges in directory feeds",
        "Protects students from opening fake files",
        "Maintains authentic institutional channels"
    ]),
    ("IMMUTABLE AUDITING", EMERALD_GREEN, [
        "Automatic diary of system operations",
        "Logs User, Action, Target, IP, & Time",
        "Helps track file leakage or illegal uploads",
        "Clean query screen for administrative audits"
    ])
]

for i, (name, color, points) in enumerate(sec_features):
    add_glass_card(slide, left_positions[i], Inches(1.6), Inches(3.7), Inches(5.1))
    
    card_box = slide.shapes.add_textbox(left_positions[i] + Inches(0.15), Inches(1.7), Inches(3.4), Inches(4.9))
    ctf = card_box.text_frame
    ctf.word_wrap = True
    
    p = ctf.paragraphs[0]
    p.text = name
    p.font.name = FONT_TITLE
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = color
    
    spc = ctf.add_paragraph()
    spc.text = "—"
    spc.font.size = Pt(10)
    spc.font.color.rgb = TEXT_MUTED
    spc.space_before = Pt(4)
    
    for pt in points:
        li = ctf.add_paragraph()
        li.text = f"✦ {pt}"
        li.font.name = FONT_BODY
        li.font.size = Pt(13)
        li.font.color.rgb = TEXT_WHITE
        li.space_before = Pt(12)

# ==================== SLIDE 9: KEY BENEFITS ====================
slide = prs.slides.add_slide(slide_layout)
apply_bg(slide)
add_title(slide, "Key Benefits to the Institution")

# Left Column: Student & Faculty Experience
add_glass_card(slide, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2))
user_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(5.2), Inches(5.0))
utf = user_box.text_frame
utf.word_wrap = True

up = utf.paragraphs[0]
up.text = "STUDENT & FACULTY EXPERIENCE"
up.font.name = FONT_TITLE
up.font.size = Pt(18)
up.font.bold = True
up.font.color.rgb = CYAN_ACCENT

benefits_user = [
    ("Premium Glassmorphism Interface", "Modern visual aesthetics built with dark theme and blur elements."),
    ("Universal Responsive Screen Compatibility", "Flawless usability across mobile phones, tablets, and computers."),
    ("Organized Resource Access", "Eliminates confusion by having files sorted in correct drawers (CSE, MCA, ECE)."),
    ("Collaborative Helpdesk", "Direct integrated ticket screen to request staff assistance.")
]
for title, desc in benefits_user:
    p_t = utf.add_paragraph()
    p_t.text = f"• {title}"
    p_t.font.name = FONT_TITLE
    p_t.font.size = Pt(14)
    p_t.font.bold = True
    p_t.font.color.rgb = TEXT_WHITE
    p_t.space_before = Pt(8)
    
    p_d = utf.add_paragraph()
    p_d.text = f"  {desc}"
    p_d.font.name = FONT_BODY
    p_d.font.size = Pt(12.5)
    p_d.font.color.rgb = TEXT_MUTED

# Right Column: Administrative Value
add_glass_card(slide, Inches(6.8), Inches(1.5), Inches(5.6), Inches(5.2))
admin_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.2), Inches(5.0))
atf2 = admin_box.text_frame
atf2.word_wrap = True

ap = atf2.paragraphs[0]
ap.text = "ADMINISTRATIVE & SAFETY VALUE"
ap.font.name = FONT_TITLE
ap.font.size = Pt(18)
ap.font.bold = True
ap.font.color.rgb = PURPLE_ACCENT

benefits_admin = [
    ("Zero Execute Security Risks", "Physical server files cannot run executable scripts, maintaining sandbox safety."),
    ("Strict Identity Registration Gate", "User registration validation binds emails to authentic institutional users."),
    ("Dynamic Audit-Driven Monitoring", "Admins keep track of storage growth and delete/audit log details instantly."),
    ("Automated Time Cleanups", "Backend sweeps database regularly to wipe expired sharing tokens.")
]
for title, desc in benefits_admin:
    s_t = atf2.add_paragraph()
    s_t.text = f"✔ {title}"
    s_t.font.name = FONT_TITLE
    s_t.font.size = Pt(14)
    s_t.font.bold = True
    s_t.font.color.rgb = EMERALD_GREEN
    s_t.space_before = Pt(8)
    
    s_d = atf2.add_paragraph()
    s_d.text = f"  {desc}"
    s_d.font.name = FONT_BODY
    s_d.font.size = Pt(12.5)
    s_d.font.color.rgb = TEXT_MUTED

# ==================== SLIDE 10: RESULTS & CONCLUSION ====================
slide = prs.slides.add_slide(slide_layout)
apply_bg(slide)
add_title(slide, "Conclusion & Review Summary")

# 2 Columns (Key Highlights & Thank You)
add_glass_card(slide, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2))
hl_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(5.2), Inches(5.0))
htf = hl_box.text_frame
htf.word_wrap = True

hp = htf.paragraphs[0]
hp.text = "PROJECT HIGHLIGHTS"
hp.font.name = FONT_TITLE
hp.font.size = Pt(18)
hp.font.bold = True
hp.font.color.rgb = CYAN_ACCENT

highlights = [
    ("Centralized Digital Drive System", "Successfully replaces disorganized email and chat file share threads."),
    ("Filename-Based Lexical Sorter", "Saves student/teacher time by sorting file categories instantly on upload."),
    ("Tokenized Deadlined Dropboxes", "Solves coursework collections securely while preserving privacy."),
    ("Comprehensive Security Controls", "Provides defense-in-depth including sanitization and obfuscation.")
]
for title, desc in highlights:
    pt = htf.add_paragraph()
    pt.text = f"✔ {title}"
    pt.font.name = FONT_TITLE
    pt.font.size = Pt(14)
    pt.font.bold = True
    pt.font.color.rgb = EMERALD_GREEN
    pt.space_before = Pt(8)
    
    pd = htf.add_paragraph()
    pd.text = desc
    pd.font.name = FONT_BODY
    pd.font.size = Pt(12.5)
    pd.font.color.rgb = TEXT_MUTED

add_glass_card(slide, Inches(6.8), Inches(1.5), Inches(5.6), Inches(5.2))
qa_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.2), Inches(5.0))
qtf = qa_box.text_frame
qtf.word_wrap = True

qp = qtf.paragraphs[0]
qp.text = "THANK YOU & QUESTIONS"
qp.font.name = FONT_TITLE
qp.font.size = Pt(24)
qp.font.bold = True
qp.font.color.rgb = PURPLE_ACCENT
qp.alignment = PP_ALIGN.CENTER
qp.space_before = Pt(15)

qp2 = qtf.add_paragraph()
qp2.text = "SECURE CAMPUS FILE SHARING & DOCUMENT MANAGEMENT SYSTEM"
qp2.font.name = FONT_TITLE
qp2.font.size = Pt(15)
qp2.font.bold = True
qp2.font.color.rgb = TEXT_WHITE
qp2.alignment = PP_ALIGN.CENTER
qp2.space_before = Pt(30)

qp3 = qtf.add_paragraph()
qp3.text = "Guided by:"
qp3.font.name = FONT_BODY
qp3.font.size = Pt(13)
qp3.font.color.rgb = TEXT_MUTED
qp3.alignment = PP_ALIGN.CENTER
qp3.space_before = Pt(15)

qp4 = qtf.add_paragraph()
qp4.text = "Dr. K. KARTHIKEYAN, MSc., MPhil., PhD., SET.,"
qp4.font.name = FONT_TITLE
qp4.font.size = Pt(14)
qp4.font.bold = True
qp4.font.color.rgb = CYAN_ACCENT
qp4.alignment = PP_ALIGN.CENTER

qp5 = qtf.add_paragraph()
qp5.text = "Assistant Professor of Computer Science\nPuratchi Thalaivi Amma Government Arts and Science College, Palladam - 641664."
qp5.font.name = FONT_BODY
qp5.font.size = Pt(12)
qp5.font.color.rgb = TEXT_MUTED
qp5.alignment = PP_ALIGN.CENTER

qp6 = qtf.add_paragraph()
qp6.text = "Open to Queries & Feedback"
qp6.font.name = FONT_TITLE
qp6.font.size = Pt(16)
qp6.font.bold = True
qp6.font.color.rgb = EMERALD_GREEN
qp6.alignment = PP_ALIGN.CENTER
qp6.space_before = Pt(25)

# Save presentation
output_filename = "d:/File-Sharing/Secure_Campus_File_Sharing_Presentation_Simple.pptx"
prs.save(output_filename)
print(f"Presentation saved successfully to: {output_filename}")

if __name__ == "__main__":
    pass
